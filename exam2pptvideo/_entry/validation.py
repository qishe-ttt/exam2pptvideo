from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import os

import click
@click.command()
@click.option("--pptx", prompt="pptx file", help="Specify the pptx file to be validated")
def validate(pptx):
  prs = Presentation(pptx)

  templates = prs.slide_layouts
  for index, t in enumerate(templates):
    slide = prs.slides.add_slide(t)
    holders = slide.shapes.placeholders

    print('Template {}-{}'.format(index+1, t.name))
    for e in holders:
      print('  %d %s' % (e.placeholder_format.idx, e.name))
   
