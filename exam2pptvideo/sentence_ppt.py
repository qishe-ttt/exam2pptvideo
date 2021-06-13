from pptx import Presentation
from exam2pptvideo.lib import readCSV, re_substitute
import json
import re

class SentencePPT:
  """It is designed as an abstract class to be inheritated for exam ppt generation of different languages,

  Note:
    **DON'T** use this class directly. To define a SentencePPT for a lanugage

      1. Define subclass inheriting ``SentencePPT``, e.g, ``ChineseSentencePPT`` 
      2. Define class variable ``_templates``, which is a :obj:`dict`, key is template genre, value is template path
      3. Define class variable ``content_keys``, which is a :obj:`list`, containing the heads in csv file
  """

  def __init__(self, sourcefile, title="", genre="classic"):
    """Initialize ppt object, ppt title, and content from source csv file

    Args:
      sourcefile (str): csv file, whose content is written into ppt
      title (str): title written in ppt home slide
      genre (str): ppt template style 
    """

    if self.__class__.__name__ != "SentencePPT":
      self._assert_class_variables()
      self._template = self.__class__._templates[genre]
      self._prs = Presentation(self._template)
      self._title = title
      self._sourcefile = sourcefile
      self._assert_content()
    else:
      raise TypeError(self.__class__.__doc__)

  def _assert_class_variables(self):
    cls = self.__class__
    assert cls._templates != None
    assert isinstance(cls._templates, dict)
    assert cls._score_code != None
    assert isinstance(cls._score_code, dict)
    assert cls._cal_slogan != None
    assert isinstance(cls._cal_slogan, str)
    assert cls.content_keys != None
    assert isinstance(cls.content_keys, list)

  def _assert_content(self):
    """Ensure the csv file content has keys defined in ppt class
    """

    content = readCSV(self._sourcefile) 
    for e in content:
      keys = e.keys()
      assert(len(keys) == len(self.__class__.content_keys))
      for k in keys:
        assert(k in self.__class__.content_keys)
    self.content = content

  def _create_opening(self):
    """Create home slide
    """

    layout = self._prs.slide_layouts.get_by_name("Cover")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "Cover"
    holders = slide.shapes.placeholders
 
    title = holders[10] 
    title.text_frame.text = self._title

  def _create_score(self, score):
    """Create score slide
    """

    layout = self._prs.slide_layouts.get_by_name("Score")
    slide = self._prs.slides.add_slide(layout)
    slide.name = score
    holders = slide.shapes.placeholders
 
    _code = self.__class__._score_code

    title, title_cn = holders[10], holders[11]
    title.text_frame.text = _code[score][0] 
    title_cn.text_frame.text = _code[score][1]

    holders[12].text_frame.text = score



  def _create_question(self, i, line):
    layout = self._prs.slide_layouts.get_by_name("Sentence")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "Q" + str(i)
    holders = slide.shapes.placeholders

    seq, question, level_checkpoint = holders[10], holders[11], holders[12]
    seq.text_frame.text = "Q." + str(i)
    question.text_frame.text = line["Question"]
    level_checkpoint.text_frame.text = line["Level"] + "-" + line["Checkpoint"]
 
  def _create_question_with_choices(self, i, line):
    layout = self._prs.slide_layouts.get_by_name("Sentence with choices")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "QWC" + str(i)
    holders = slide.shapes.placeholders

    seq, question = holders[10], holders[11]
    seq.text_frame.text = "Q." + str(i)
    question.text_frame.text = line["Question"]

    A, B, C, D = holders[12], holders[13], holders[14], holders[15]
    A.text_frame.text = line["A"]
    B.text_frame.text = line["B"]
    C.text_frame.text = line["C"]
    D.text_frame.text = line["D"]

    level_checkpoint = holders[16]
    level_checkpoint.text_frame.text = line["Level"] + "-" + line["Checkpoint"]
 
  def _create_answer(self, i, line):
    layout = self._prs.slide_layouts.get_by_name("Sentence analysis")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "A" + str(i)
    holders = slide.shapes.placeholders

    seq, correct, question = holders[10], holders[11], holders[12]

    seq.text_frame.text = "A." + str(i)
    correct.text_frame.text = line["Correct"]
    
    replaced_question = re_substitute('_+', line["Question"], line[line["Correct"]])

    question.text_frame.text = replaced_question

    level, checkpoint, explanation = holders[13], holders[14], holders[15]
    level.text_frame.text = line["Level"]
    checkpoint.text_frame.text = line["Checkpoint"]
    explanation.text_frame.text = line["Explanation"]

  def _create_answer_with_choices(self, i, line):
    layout = self._prs.slide_layouts.get_by_name("Sentence analysis with choices")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "AWC" + str(i)
    holders = slide.shapes.placeholders

    seq, correct, question = holders[10], holders[11], holders[12]

    seq.text_frame.text = "A." + str(i)
    correct.text_frame.text = line["Correct"]

    replaced_question = re_substitute('_+', line["Question"], line[line["Correct"]])

    question.text_frame.text = replaced_question

    A, B, C, D = holders[13], holders[14], holders[15], holders[16] 
    A.text_frame.text = line["A"]
    B.text_frame.text = line["B"]
    C.text_frame.text = line["C"]
    D.text_frame.text = line["D"]

    level, checkpoint, explanation = holders[17], holders[18], holders[19]
    level.text_frame.text = line["Level"]
    checkpoint.text_frame.text = line["Checkpoint"]
    explanation.text_frame.text = line["Explanation"]

  def _create_calculation(self):
    layout = self._prs.slide_layouts.get_by_name("Calculation")
    slide = self._prs.slides.add_slide(layout)
    holders = slide.shapes.placeholders
    
    cal = holders[10]
    _cal_slogan = self.__class__._cal_slogan

    cal.text_frame.text = _cal_slogan 


    slide.name = "Calculation"

  def _create_wrong(self):
    layout = self._prs.slide_layouts.get_by_name("Wrong")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "Wrong"

  def _create_correct(self):
    layout = self._prs.slide_layouts.get_by_name("Correct")
    slide = self._prs.slides.add_slide(layout)
    slide.name = "Correct"

  def _create_QAs(self):
    for i, line in enumerate(self.content):
      self._create_question(i+1, line)
      self._create_answer(i+1, line)
      self._create_question_with_choices(i+1, line)
      self._create_answer_with_choices(i+1, line)
    
  def _save_ppt(self, destfile):
    """Save ppt object into file
    """

    self._prs.save(destfile)

  def convert_to_ppt(self, destfile='test.pptx'):
    """Convert csv file containing exam information into pptx file

    Args:
      destfile (str): pptx file path
    """

    self._create_opening()
    self._create_QAs()
    self._create_wrong()
    self._create_correct()
    self._create_calculation()
    self._create_score("100")
    self._create_score("60-90")
    self._create_score("0-50")  

    self._save_ppt(destfile)    

