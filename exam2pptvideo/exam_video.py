from pptx import Presentation 
from exam2pptvideo.lib import pptx2pdf, pdf2images 
from subprocess import call, DEVNULL
import hashlib
import os
from abc import abstractmethod

class ExamVideo:
  """It is designed as an abstract class to be inheritated for exam video generation of different languages,

  Note:
    **DON'T** use this class directly. To define a ExamVideo for a lanugage

      1. Define subclass inheriting ``ExamVideo``, e.g, ``ChineseExamVideo`` 
      2. Define class variable ``_soundindex``, which is a :obj:`dict`, key is slide type, value is sound path
  """

  def __init__(self, sourceppt):
    """Initialize video object, slide from source pptx file

    Args:
      sourceppt (str): pptx file
    """


    if self.__class__.__name__ != "ExamVideo":
      self._assert_class_variables()
      self._sourceppt = sourceppt
      self.slides = Presentation(sourceppt).slides
      _code =  hashlib.sha224(sourceppt.encode('utf-8')).hexdigest()[:8]
      self._img_folder = _code + "_imgs" 
    else:
      raise TypeError(self.__class__.__doc__)

  def _assert_class_variables(self):
    cls = self.__class__
    assert cls._soundindex != None
    assert isinstance(cls._soundindex, dict)

  def _create_jpgs(self):

    if not os.path.isdir(self._img_folder):
      os.mkdir(self._img_folder)

    pdf = pptx2pdf(self._sourceppt, self._img_folder)
    pdf2images(pdf, self._img_folder)

  def create_videos(self, video_folder):
    """Generate videos from pptx

    Args:
      video_folder (str): path to store videos

    """

    self._create_jpgs()

    if not os.path.isdir(video_folder):
      os.mkdir(video_folder) 

    size = len(self.slides)
    for i in range(size):
      slide_name = self.slides[i].name
      m4a = self._get_slide_sound(i)
      jpg = "{}/{}.jpg".format(self._img_folder, i)
      mp4 = "{}/{}.mp4".format(video_folder, slide_name)
     
      if not os.path.isfile(mp4):
        call(["ffmpeg", "-loop", "1", "-i", jpg, "-i", m4a, "-vf", "scale=1920:1080", "-c:v", "libx264", "-c:a", "copy", "-shortest", "-movflags", "+faststart", mp4], stdout=DEVNULL)
        #call(["ffmpeg", "-loop", "1", "-i", jpg, "-i", mp3, "-vf", "scale=1920:1080", "-c:v", "libx264", "-tune", "stillimage", "-c:a", "aac", "-pix_fmt", "yuv420p", "-shortest", mp4], stdout=DEVNULL)
 
  @abstractmethod
  def _get_slide_sound(self, i):
    """Create sound for slide
    """
    pass

